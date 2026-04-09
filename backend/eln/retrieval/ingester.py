"""
Document ingester for PDF, PPTX, and CSV files.

Converts raw files into LangChain Document objects suitable for
chunking and embedding. Each Document carries metadata for downstream
citation construction.
"""

import csv
import hashlib
from pathlib import Path

import pdfplumber
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from docx import Document as DocxDocument


class IngestionError(Exception):
    """Raised when a file cannot be parsed (corrupt, empty, unsupported)."""

# Chunking tuned for scientific SOPs: 512 tokens is large enough to keep
# protocol context together, small enough to score precisely in retrieval.
_SPLITTER = RecursiveCharacterTextSplitter(
    chunk_size=512,
    chunk_overlap=64,
    length_function=len,  # character-based; close enough for English text
    separators=["\n\n", "\n", ". ", " ", ""],
)


def _sha256(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def _doc_id(file_hash: str, chunk_index: int) -> str:
    """Deterministic document ID: same file + same chunk index = same ID."""
    return hashlib.sha256(f"{file_hash}:{chunk_index}".encode()).hexdigest()[:16]


class Ingester:
    """
    Parses files into chunked LangChain Documents with metadata.

    Supported formats: .pdf, .pptx, .csv
    Each Document.metadata includes:
      - source_path, sha256, doc_id, chunk_index, total_chunks
      - page (for PDF), slide (for PPTX)
    """

    def ingest(
        self,
        path: Path,
        source_type: str = "paper",
        collection_name: str = "papers",
        user_id: str = "default",
        storage_backend=None,
    ) -> list[Document]:
        """Route to the appropriate parser based on file extension.

        If storage_backend is provided, the source file is stored there
        (content-addressed by sha256) and an s3_key is added to each chunk's
        metadata so citations can link back to the original file.

        Raises IngestionError on corrupt or unreadable files.
        """
        suffix = path.suffix.lower()
        try:
            if suffix == ".pdf":
                chunks = self.ingest_pdf(path, source_type)
            elif suffix == ".pptx":
                chunks = self.ingest_pptx(path, source_type)
            elif suffix == ".csv":
                chunks = self.ingest_csv(path, source_type)
            else:
                raise ValueError(f"Unsupported file type: {suffix}")
        except (ValueError, FileNotFoundError):
            raise
        except Exception as e:
            raise IngestionError(f"Failed to ingest {path.name}: {e}") from e

        if storage_backend is not None and chunks:
            file_hash = chunks[0].metadata.get("sha256") or _sha256(path)
            key = f"kb/{user_id}/{collection_name}/{file_hash}{suffix}"
            if not storage_backend.exists(key):
                storage_backend.put_object(key, path)
            for chunk in chunks:
                chunk.metadata["s3_key"] = key

        return chunks

    def ingest_pdf(self, path: Path, source_type: str = "paper") -> list[Document]:
        """
        Extract text from PDF using pdfplumber (handles tables well).
        Returns one Document per chunk, with page-level metadata.
        """
        file_hash = _sha256(path)
        page_docs: list[Document] = []

        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""

                # Also extract tables as pipe-separated text
                for table in page.extract_tables():
                    if table:
                        rows = [" | ".join(str(cell or "") for cell in row) for row in table]
                        text += "\n" + "\n".join(rows) + "\n"

                if text.strip():
                    page_docs.append(Document(
                        page_content=text.strip(),
                        metadata={
                            "source_path": str(path),
                            "sha256": file_hash,
                            "source_type": source_type,
                            "page": page_num,
                        },
                    ))

        return self._chunk_and_tag(page_docs, file_hash, source_type, str(path))

    def ingest_pptx(self, path: Path, source_type: str = "paper") -> list[Document]:
        """Extract text from PowerPoint slides."""
        file_hash = _sha256(path)
        slide_docs: list[Document] = []

        prs = Presentation(str(path))
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
            if parts:
                slide_docs.append(Document(
                    page_content="\n".join(parts),
                    metadata={
                        "source_path": str(path),
                        "sha256": file_hash,
                        "source_type": source_type,
                        "slide": slide_num,
                    },
                ))

        return self._chunk_and_tag(slide_docs, file_hash, source_type, str(path))

    def ingest_csv(self, path: Path, source_type: str = "paper") -> list[Document]:
        """
        Convert CSV to text: column names + first N rows + summary stats.
        Treats the entire CSV as one document (then chunks it).
        """
        file_hash = _sha256(path)
        max_rows = 100

        with open(path, newline="", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = []
            for i, row in enumerate(reader):
                rows.append(row)
                if i >= max_rows:
                    break

        if not rows:
            return []

        header = rows[0]
        data_rows = rows[1:]

        lines = [
            f"CSV file: {path.name}",
            f"Columns ({len(header)}): {', '.join(header)}",
            f"Rows shown: {len(data_rows)} (of file)",
            "",
        ]
        for row in data_rows:
            lines.append(" | ".join(row))

        text = "\n".join(lines)
        page_docs = [Document(
            page_content=text,
            metadata={
                "source_path": str(path),
                "sha256": file_hash,
                "source_type": source_type,
            },
        )]

        return self._chunk_and_tag(page_docs, file_hash, source_type, str(path))

    def extract_full_text(self, path: Path) -> str:
        """Return all text from a document as a single string (non-chunked).

        Used for template extraction where structure analysis needs the full
        document rather than isolated chunks.

        Supports: .pdf, .pptx, .docx, .txt, .md
        """
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return self._full_text_pdf(path)
        if suffix == ".pptx":
            return self._full_text_pptx(path)
        if suffix == ".docx":
            return self._full_text_docx(path)
        if suffix in (".txt", ".md"):
            return path.read_text(errors="replace")
        raise IngestionError(f"Unsupported format for full-text extraction: {suffix}")

    def _full_text_pdf(self, path: Path) -> str:
        parts: list[str] = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                for table in page.extract_tables():
                    if table:
                        rows = [" | ".join(str(cell or "") for cell in row) for row in table]
                        text += "\n" + "\n".join(rows)
                if text.strip():
                    parts.append(text.strip())
        return "\n\n".join(parts)

    def _full_text_pptx(self, path: Path) -> str:
        parts: list[str] = []
        prs = Presentation(str(path))
        for slide in prs.slides:
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
            if slide_parts:
                parts.append("\n".join(slide_parts))
        return "\n\n".join(parts)

    def _full_text_docx(self, path: Path) -> str:
        doc = DocxDocument(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _chunk_and_tag(
        self,
        docs: list[Document],
        file_hash: str,
        source_type: str,
        source_path: str,
    ) -> list[Document]:
        """Split documents into chunks and assign deterministic doc_ids."""
        chunks = _SPLITTER.split_documents(docs)
        total = len(chunks)
        for i, chunk in enumerate(chunks):
            chunk.metadata["doc_id"] = _doc_id(file_hash, i)
            chunk.metadata["chunk_index"] = i
            chunk.metadata["total_chunks"] = total
            chunk.metadata.setdefault("source_path", source_path)
            chunk.metadata.setdefault("sha256", file_hash)
            chunk.metadata.setdefault("source_type", source_type)
        return chunks
